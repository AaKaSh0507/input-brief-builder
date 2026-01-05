import os
import shutil
from typing import Optional
from pathlib import Path
import mimetypes
import logging
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
import csv
from pptx import Presentation
from PIL import Image
import io

logger = logging.getLogger(__name__)

class DocumentService:
    def __init__(self):
        self.upload_dir = Path("/app/backend/uploads")
        self.upload_dir.mkdir(exist_ok=True)
    
    async def save_file(self, file_content: bytes, filename: str) -> tuple[str, str]:
        """Save uploaded file and return file path and mime type"""
        try:
            file_path = self.upload_dir / filename
            
            with open(file_path, "wb") as f:
                f.write(file_content)
            
            mime_type = mimetypes.guess_type(filename)[0] or "application/octet-stream"
            
            return str(file_path), mime_type
        except Exception as e:
            logger.error(f"Error saving file: {str(e)}")
            raise
    
    async def extract_text_content(self, file_path: str, file_type: str) -> Optional[str]:
        """Extract text content from various file types"""
        try:
            if file_type == "pdf" or file_path.endswith(".pdf"):
                return await self._extract_from_pdf(file_path)
            elif file_type in ["doc", "docx"] or file_path.endswith((".doc", ".docx")):
                return await self._extract_from_docx(file_path)
            elif file_type in ["xlsx", "xls"] or file_path.endswith((".xlsx", ".xls")):
                return await self._extract_from_excel(file_path)
            elif file_type == "csv" or file_path.endswith(".csv"):
                return await self._extract_from_csv(file_path)
            elif file_type in ["ppt", "pptx"] or file_path.endswith((".ppt", ".pptx")):
                return await self._extract_from_pptx(file_path)
            elif file_type in ["jpg", "jpeg", "png", "gif"] or file_path.endswith((".jpg", ".jpeg", ".png", ".gif")):
                return "[Image file - use AI for analysis]"
            else:
                return None
        except Exception as e:
            logger.error(f"Error extracting text: {str(e)}")
            return None
    
    async def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF"""
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error reading PDF: {str(e)}")
            return ""
    
    async def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX"""
        try:
            doc = DocxDocument(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text.strip()
        except Exception as e:
            logger.error(f"Error reading DOCX: {str(e)}")
            return ""
    
    async def _extract_from_csv(self, file_path: str) -> str:
        """Extract text from CSV"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                rows = list(reader)
                text = "\n".join([", ".join(row) for row in rows])
                return text
        except Exception as e:
            logger.error(f"Error reading CSV: {str(e)}")
            return ""
    
    async def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error reading PPTX: {str(e)}")
            return ""
    
    def delete_file(self, file_path: str):
        """Delete a file"""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception as e:
            logger.error(f"Error deleting file: {str(e)}")

# Singleton instance
document_service = DocumentService()
